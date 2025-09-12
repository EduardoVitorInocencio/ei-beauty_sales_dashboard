import pytest
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from app.services.template_service import update_template
import os

def create_mock_template_with_charts(path):
    """
    Cria um template PPTX com 3 slides e gráficos nomeados conforme padrão.
    """
    ppt = Presentation()

    # Slide 1 - Produto
    slide1 = ppt.slides.add_slide(ppt.slide_layouts[5])
    chart_data1 = CategoryChartData()
    chart_data1.categories = ["Placeholder"]
    chart_data1.add_series("Vendas", [0])
    chart1 = slide1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(5), Inches(3), chart_data1
    ).chart
    slide1.shapes[-1].name = "Chart_Produto"

    # Slide 2 - Região
    slide2 = ppt.slides.add_slide(ppt.slide_layouts[5])
    chart_data2 = CategoryChartData()
    chart_data2.categories = ["Placeholder"]
    chart_data2.add_series("Vendas", [0])
    chart2 = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(5), Inches(3), chart_data2
    ).chart
    slide2.shapes[-1].name = "Chart_Regiao"

    # Slide 3 - Mensal
    slide3 = ppt.slides.add_slide(ppt.slide_layouts[5])
    chart_data3 = CategoryChartData()
    chart_data3.categories = ["Placeholder"]
    chart_data3.add_series("Vendas", [0])
    chart3 = slide3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(5), Inches(3), chart_data3
    ).chart
    slide3.shapes[-1].name = "Chart_Mensal"

    ppt.save(path)

def create_mock_excel(path):
    """
    Cria um Excel temporário com dados fictícios de vendas.
    """
    df = pd.DataFrame({
        "produto": ["A", "B"],
        "regiao": ["Norte", "Sul"],
        "data": pd.to_datetime(["2025-01-01", "2025-02-01"]),
        "vendas": [100, 200]
    })
    df.to_excel(path, index=False)

def test_update_template_full(tmp_path):
    """
    Testa o fluxo completo de atualização do template.
    """
    excel_path = tmp_path / "vendas.xlsx"
    template_path = tmp_path / "template.pptx"
    output_path = tmp_path / "template_atualizado.pptx"

    create_mock_excel(excel_path)
    create_mock_template_with_charts(template_path)

    # Rodar o serviço
    updated_file = update_template(str(excel_path), str(template_path), str(output_path))

    # Verifica se o arquivo foi criado
    assert os.path.exists(updated_file)

    # Verifica se os gráficos foram atualizados com dados corretos
    ppt = Presentation(updated_file)

    # Slide 1 - Produto
    slide1 = ppt.slides[0]
    chart1 = [s for s in slide1.shapes if s.has_chart and s.name == "Chart_Produto"][0].chart
    series1 = chart1.series[0].values
    assert sum(series1) == 300  # 100 + 200

    # Slide 2 - Região
    slide2 = ppt.slides[1]
    chart2 = [s for s in slide2.shapes if s.has_chart and s.name == "Chart_Regiao"][0].chart
    series2 = chart2.series[0].values
    assert sum(series2) == 300

    # Slide 3 - Mensal
    slide3 = ppt.slides[2]
    chart3 = [s for s in slide3.shapes if s.has_chart and s.name == "Chart_Mensal"][0].chart
    series3 = chart3.series[0].values
    assert sum(series3) == 300
