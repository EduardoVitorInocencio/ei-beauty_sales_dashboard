from pptx.chart.data import CategoryChartData
from pptx import Presentation
from app.graphs.base_graph import BaseGraph
from app.data.load_data_from_sql import load_data_from_db
from pptx.dml.color import RGBColor
from pptx.util import Pt
import os
import json

class ProductSalesGraph(BaseGraph):
    """
    Gráfico de vendas por produto (Página 1).
    Atualiza o gráfico existente no template pelo nome padronizado.
    """

    CHART_NAME = "Chart_Produto"  # Nome do gráfico no template

    # Filtros fixos
    FILTERS = {
        "formats": "T. Embalagens",
        "region": "Total Brasil",
        "channels": "Hipermercados",
        "periods": "YTD Dec-20",
        "brands": "Coca-Cola Cia"
    }

    # Colunas que queremos no DataFrame final
    SELECT_COLS = [
        "periods",
        "units","value_(r$)","weighted_hholds","weighted_buyers","penetration_(%)",
        "units_per_buyer","spend_per_buyer_(r$)","units_per_trip","spend_per_trip",
        "frequency","avg_price_per_unit"
    ]

    COLS_RENAME = {
        "units": "Units",
        "value_(r$)": "Value_(R$)",
        "weighted_hholds": "Weighted_HHOLDS",
        "weighted_buyers": "Weighted_Buyers",
        "penetration_(%)": "Penetration_(%)",
        "units_per_buyer": "Units_per_Buyer",
        "spend_per_buyer_(r$)": "Spend_per_Buyer_(R$)",
        "units_per_trip": "Units_per_Trip",
        "spend_per_trip": "Spend_per_Trip",
        "frequency": "Frequency",
        "avg_price_per_unit": "Avg_Price_per_Unit"
    }

    def __init__(self, ppt: Presentation):
        # Busca os dados direto do SQLite
        data = load_data_from_db()
        super().__init__(data, ppt)


    def filter_data(self):
        """
        Aplicar os filtros e selecionar/renomear as colunas desejadas.
        """
        df = self.data.copy()
        for col, val in self.FILTERS.items():
            df = df[df[col] == val]
        # Seleciona e renomeia colunas
        self.filtered_data = df[self.SELECT_COLS].rename(columns=self.COLS_RENAME)

    def update_chart(self):
        # Implementação mínima para satisfazer a classe abstrata
        pass

    def update_value_shape(self):
        """
        Atualiza todos os shapes 'abs' usando os valores do DataFrame já filtrado e renomeado.
        Mantém fonte tamanho 8, cor preta e Calibri.
        """
        import os
        import json
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        # Carrega JSON
        json_path = os.path.join(os.path.dirname(__file__), "..", "data", "shapes_mapping.json")
        with open(json_path, "r", encoding="utf-8") as f:
            shapes_map = json.load(f)

        abs_shapes = shapes_map.get("abs", [])
        if not abs_shapes:
            print("Nenhum shape 'abs' encontrado no JSON.")
            return

        # Agora aponta para as colunas já renomeadas
        shape_to_col = {
            "shape_abs_volume_units": "Units",
            "shape_abs_valor": "Value_(R$)",
            "shape_abs_total_domicilios": "Weighted_HHOLDS",
            "shape_abs_compradores": "Weighted_Buyers",
            "shape_absr_per_penetracao": "Penetration_(%)",
            "shape_abs_volume_por_comprador": "Units_per_Buyer",
            "shape_abs_gasto_por_comprador": "Spend_per_Buyer_(R$)",
            "shape_abs_volume_por_viagem": "Units_per_Trip",
            "shape_abs_gasto_por_viagem": "Spend_per_Trip",
            "shape_abs_frequencia": "Frequency",
            "shape_abs_frequencia2": "Frequency",
            "shape_abs_preco_medio": "Avg_Price_per_Unit"
        }

        def find_and_update(shape, target_name, value):
            if shape.has_text_frame and shape.name.strip().lower() == target_name.lower():
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(value)
                run.font.size = Pt(8)
                run.font.color.rgb = RGBColor(0, 0, 0)
                run.font.name = "Calibri"
                print(f"Shape '{shape.name}' atualizado com {value}")
                return True

            if shape.shape_type == 6:  # grupo
                for subshape in shape.shapes:
                    if find_and_update(subshape, target_name, value):
                        return True
            return False

        # Atualiza cada shape
        for shape_name in abs_shapes:
            df_col = shape_to_col.get(shape_name.lower())
            if df_col and df_col in self.filtered_data.columns:
                value = self.filtered_data[df_col].sum()
            else:
                value = "N/A"

            updated = False
            for slide in self.ppt.slides:
                for shape in slide.shapes:
                    if find_and_update(shape, shape_name, value):
                        updated = True
                        break
                if updated:
                    break

            if not updated:
                print(f"Shape '{shape_name}' não encontrado no template.")
        
